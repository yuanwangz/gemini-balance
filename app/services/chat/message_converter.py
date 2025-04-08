# app/services/chat/message_converter.py

from abc import ABC, abstractmethod
import re
from typing import Any, Dict, List, Optional
import requests
import base64
import json
SUPPORTED_ROLES = ["user", "model", "system"]
IMAGE_URL_PATTERN = r'\[image\]\((.*?)\)'


class MessageConverter(ABC):
    """消息转换器基类"""

    @abstractmethod
    def convert(self, messages: List[Dict[str, Any]]) -> tuple[List[Dict[str, Any]], Optional[Dict[str, Any]]]:
        pass

def _get_mime_type_and_data(base64_string):
    """
    从 base64 字符串中提取 MIME 类型和数据。
    
    参数:
        base64_string (str): 可能包含 MIME 类型信息的 base64 字符串
        
    返回:
        tuple: (mime_type, encoded_data)
    """
    # 检查字符串是否以 "data:" 格式开始
    if base64_string.startswith('data:'):
        # 提取 MIME 类型和数据
        pattern = r'data:([^;]+);base64,(.+)'
        match = re.match(pattern, base64_string)
        if match:
            mime_type = "image/jpeg" if match.group(1) == "image/jpg" else match.group(1)
            encoded_data = match.group(2)
            return mime_type, encoded_data
    
    # 如果不是预期格式，假定它只是数据部分
    return None, base64_string

def _convert_file(file_url: str, api_key: str) -> Dict[str, Any]:
    mime_type = None
    encoded_data = None
    file_size = None
    
    original_url = file_url  # 保存原始URL
    
    if file_url.startswith("http"):
        encoded_data, mime_type = _convert_file_to_base64(file_url)
        # 不再覆盖file_url变量
    elif file_url.startswith("data:"):  # 使用elif而不是if
        mime_type, encoded_data = _get_mime_type_and_data(file_url)
    
    print(f"mime_type: {mime_type}")
    
    if _is_office_document(mime_type):
        encoded_data, mime_type = convert_to_text(encoded_data, mime_type)
    
    file_size = calculate_image_size(encoded_data)
    print(f"file_size: {file_size} MB")
    
    if file_size < 20:
        return {
            "inline_data": {
                "mime_type": mime_type,
                "data": encoded_data
            }
        }
    else:
        # 使用original_url或encoded_data上传
        file_uri, upload_mime_type = upload_file_to_gemini(original_url if original_url.startswith("http") else encoded_data, api_key, mime_type)
        return {
            "file_data": {
                "mime_type": upload_mime_type,
                "file_uri": file_uri
            }
        }


def _convert_file_to_base64(url: str) -> tuple[str, str]:
    """
    将文件URL转换为base64编码
    Args:
        url: 文件URL
    Returns:
        tuple: (base64编码的文件数据, MIME类型)
    """
    response = requests.get(url)
    if response.status_code == 200:
        # 将文件内容转换为base64
        file_data = base64.b64encode(response.content).decode('utf-8')
        mime_type = response.headers.get('Content-Type', 'application/octet-stream')
        
        # 如果是通用二进制类型，尝试根据URL或文件头检测实际类型
        if mime_type == 'application/octet-stream':
            # 1. 从URL检测
            import os
            file_ext = os.path.splitext(url.split('?')[0].lower())[1]
            
            # 添加对常见图片格式的支持
            if file_ext in ['.png']:
                mime_type = 'image/png'
            elif file_ext in ['.jpg', '.jpeg', '.jpe']:
                mime_type = 'image/jpeg'
            elif file_ext in ['.gif']:
                mime_type = 'image/gif'
            elif file_ext in ['.webp']:
                mime_type = 'image/webp'
            elif file_ext in ['.svg']:
                mime_type = 'image/svg+xml'
            elif file_ext in ['.bmp']:
                mime_type = 'image/bmp'
            elif file_ext in ['.tiff', '.tif']:
                mime_type = 'image/tiff'
            # 原有的Office文档支持
            elif file_ext in ['.doc', '.dot']:
                mime_type = 'application/msword'
            elif file_ext in ['.docx', '.docm', '.dotx', '.dotm']:
                mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif file_ext in ['.xls', '.xlt', '.xla']:
                mime_type = 'application/vnd.ms-excel'
            elif file_ext in ['.xlsx', '.xlsm', '.xltx', '.xltm']:
                mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif file_ext in ['.ppt', '.pot', '.pps', '.ppa']:
                mime_type = 'application/vnd.ms-powerpoint'
            elif file_ext in ['.pptx', '.pptm', '.potx', '.potm', '.ppsx']:
                mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            elif file_ext in ['.html', '.htm']:
                mime_type = 'text/html'
            elif file_ext in ['.xhtml']:
                mime_type = 'application/xhtml+xml'
            
            # 2. 如果URL没有提供足够信息，可以检查文件头部字节
            if mime_type == 'application/octet-stream':
                binary_data = base64.b64decode(file_data)
                # 添加对常见图片格式的二进制头检测
                if binary_data.startswith(b'\x89PNG\r\n\x1a\n'):
                    mime_type = 'image/png'
                elif binary_data.startswith(b'\xff\xd8\xff'):
                    mime_type = 'image/jpeg'
                elif binary_data.startswith(b'GIF87a') or binary_data.startswith(b'GIF89a'):
                    mime_type = 'image/gif'
                elif binary_data.startswith(b'RIFF') and b'WEBP' in binary_data[0:16]:
                    mime_type = 'image/webp'
                # 原有的文件类型检测
                elif binary_data[:4] == b'PK\x03\x04':
                    mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                # Word .doc (旧格式)
                elif binary_data[:8] in [b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1', b'\xCF\x11\xE0\xA1\xB1\x1A\xE1\x00']:
                    mime_type = 'application/msword'
                # 检查是否为HTML文件（查找HTML标签的特征）
                elif binary_data.lower().startswith(b'<!doctype html') or binary_data.lower().startswith(b'<html'):
                    mime_type = 'text/html'
        
        # 当MIME类型未设置或为空时，基于内容推测
        if not mime_type or mime_type == 'application/octet-stream':
            # 尝试检测HTML
            try:
                content_start = response.content[:100].lower().decode('utf-8', errors='ignore')
                if '<!doctype html' in content_start or '<html' in content_start:
                    mime_type = 'text/html'
            except:
                # 解码失败或其他错误时，保持原样
                pass
                
        print(f"检测到MIME类型: {mime_type}")
        return file_data, mime_type
    else:
        raise Exception(f"Failed to fetch file: {response.status_code}")


def _process_text_with_file(text: str, api_key: str) -> List[Dict[str, Any]]:
    """
    处理可能包含文件URL的文本，提取文件并转换为base64

    Args:
        text: 可能包含文件URL的文本

    Returns:
        List[Dict[str, Any]]: 包含文本和文件的部分列表
    """
    parts = []
    img_url_match = re.search(IMAGE_URL_PATTERN, text)
    if img_url_match:
        # 提取URL
        img_url = img_url_match.group(1)
        
        # 检查URL是否是有效URL（不是占位符如"链接地址"）
        if img_url == "链接地址" or not (img_url.startswith("http") or img_url.startswith("data:")):
            # 占位符或无效URL，作为纯文本处理
            parts.append({"text": text})
            return parts
        
        # 将URL对应的图片转换为base64
        try:
            inline_data = _convert_file(img_url, api_key)
            parts.append(inline_data)
        except Exception as e:
            print(f"转换文件失败: {str(e)}")
            # 如果转换失败，回退到文本模式
            parts.append({"text": text})
    else:
        # 没有图片URL，作为纯文本处理
        parts.append({"text": text})
    return parts

def calculate_image_size(base64_data):
    # 解码base64数据
    image_data = base64.b64decode(base64_data)
    
    # 计算二进制数据的长度
    data_length = len(image_data)
    
    # 将长度转换为文件大小（以MB为单位）
    file_size = data_length / (1024 * 1024)  # MB
    return file_size

class OpenAIMessageConverter(MessageConverter):
    """OpenAI消息格式转换器"""

    def convert(self, messages: List[Dict[str, Any]], api_key: str) -> tuple[List[Dict[str, Any]], Optional[Dict[str, Any]]]:
        converted_messages = []
        system_instruction_parts = []

        for idx, msg in enumerate(messages):
            role = msg.get("role", "")
            if role not in SUPPORTED_ROLES:
                if role == "tool":
                    role = "function"
                if role == "assistant":
                    role = "model"
                # else:
                #     # 如果是最后一条消息，则认为是用户消息
                #     if idx == len(messages) - 1:
                #         role = "user"
                #     else:
                #         role = "model"

            parts = []
            # 特别处理最后一个assistant的消息，按\n\n分割
            if role == "model" and idx == len(messages) - 2 and msg.get("content") and isinstance(msg.get("content"), str):
                # 按\n\n分割消息
                content_parts = msg["content"].split("\n\n")
                for part in content_parts:
                    if not part.strip():  # 跳过空内容
                        continue
                    # 处理可能包含图片的文本
                    parts.extend(_process_text_with_file(part,api_key))
                
            elif role == "function":
                # 处理工具返回的消息 - Gemini格式为functionResponse
                # 先确保有name字段，如果没有则尝试使用tool_call_id
                print(f"msg: {msg}")
                content = msg.get("content")
                print(f"content: {content}")
                function_name = msg.get("name") or msg.get("tool_call_id") or "unknown_function"
                # 转换为Gemini的functionResponse格式
                parts.append({
                    "functionResponse": {
                        "name": function_name,
                        "response": {
                            "content": {
                                "result": msg.get("content")
                            }
                        }
                    }
                })
            elif isinstance(msg["content"], str) and msg["content"]:
                # 请求 gemini 接口时如果包含 content 字段但内容为空时会返回 400 错误，所以需要判断是否为空并移除
                parts.extend(_process_text_with_file(msg["content"],api_key))
            elif isinstance(msg["content"], list):
                for content in msg["content"]:
                    if isinstance(content, str) and content:
                        parts.append({"text": content})
                    elif isinstance(content, dict):
                        if content["type"] == "text" and content["text"]:
                            parts.append({"text": content["text"]})
                        elif content["type"] == "image_url":
                            parts.append(_convert_file(content["image_url"]["url"],api_key))
                            
            if role == "model":
                tool_calls = msg.get("tool_calls", [])
                if tool_calls:
                    for tool_call in tool_calls:
                        args = tool_call.get("function").get("arguments")
                        # 判断是否为json字符串,如果不是则尝试转换
                        if isinstance(args, str):
                            try:
                                args = json.loads(args)
                            except json.JSONDecodeError:
                                print(f"args is not a json string: {args}")
                                pass
                        parts.append({
                            "functionCall": {
                                "name": tool_call.get("function").get("name"),
                                "args": args
                            }
                        })
            if parts:
                if role == "system":
                    system_instruction_parts.extend(parts)
                else:
                    converted_messages.append({"role": role, "parts": parts})

        system_instruction = (
            None
            if not system_instruction_parts
            else {
                "role": "system",
                "parts": system_instruction_parts,
            }
        )
        # print(f"converted_messages: {converted_messages}")
        return converted_messages, system_instruction
    
def _is_office_document(mime_type: str) -> bool:
    """判断文件是否为Office文档"""
    office_mime_types = [
        'application/msword',  # .doc
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # .docx
        'application/vnd.ms-excel',  # .xls
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',  # .xlsx
        'application/vnd.ms-powerpoint',  # .ppt
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # .pptx
        'application/vnd.oasis.opendocument.text',  # .odt
        'application/vnd.oasis.opendocument.spreadsheet',  # .ods
        'application/vnd.oasis.opendocument.presentation'  # .odp
    ]
    return mime_type in office_mime_types
 
def convert_to_text(file_base64: str, mime_type: str) -> tuple[str, str]:
    """
    将Office文档转换为纯文本
    
    Args:
        file_base64: 文件内容的base64编码
        mime_type: 文件的MIME类型
        
    Returns:
        tuple: (文本内容的base64编码, 文本的MIME类型)
    """
    import tempfile
    import os
    import base64
    
    # 解码base64获取文件内容
    file_content = base64.b64decode(file_base64)
     
    try:
        # 确定文件类型
        extension = '.bin'
        if mime_type == 'application/msword':
            extension = '.doc'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            extension = '.docx'
        elif mime_type == 'application/vnd.ms-excel':
            extension = '.xls'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            extension = '.xlsx'
        elif mime_type == 'application/vnd.ms-powerpoint':
            extension = '.ppt'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            extension = '.pptx'
        elif mime_type == 'text/html' or mime_type == 'application/xhtml+xml':
            extension = '.html'
        with tempfile.NamedTemporaryFile(delete=False, suffix=extension) as temp_file:
            temp_file.write(file_content)
            temp_path = temp_file.name
        
        file_type = None
        if extension in ['.doc', '.docx']:
            file_type = 'word'
        elif extension in ['.xls', '.xlsx']:
            file_type = 'excel'
        elif extension in ['.ppt', '.pptx']:
            file_type = 'powerpoint'
        elif extension in ['.html', '.xhtml']:
            file_type = 'html'
        # 根据文件类型转换
        text_content = ""
        print(f"file_type: {file_type}")
        if file_type == 'html':
            # 处理HTML文件 - 读取原始内容
            try:
                # 直接读取HTML文件内容
                with open(temp_path, 'r', encoding='utf-8', errors='replace') as html_file:
                    text_content = html_file.read()
            except Exception as e:
                text_content = f"读取HTML文件失败: {str(e)}"
                print(f"HTML解析错误: {str(e)}")
        elif file_type == 'word':
            # 处理Word文档
            from docx import Document
            try:
                doc = Document(temp_path)
                # 简单提取段落文本
                paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
                print(f"找到 {len(paragraphs)} 个段落")
                
                # 处理表格内容（对合并单元格进行特殊处理）
                table_contents = []
                
                # 只有在没有找到段落文本或者段落文本很少时才考虑表格内容
                if len(paragraphs) < 5:  # 如果段落少于5个，也考虑表格内容
                    for table_idx, table in enumerate(doc.tables):
                        # 创建一个哈希表来记录已处理的文本内容
                        seen_text = set()
                        table_text = []
                        
                        # 对表格每行处理
                        for row in table.rows:
                            row_text_parts = []
                            
                            # 过滤掉重复的单元格文本
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text and cell_text not in seen_text:
                                    row_text_parts.append(cell_text)
                                    seen_text.add(cell_text)
                            
                            # 仅当有非空内容时添加该行
                            if row_text_parts:
                                row_text = " | ".join(row_text_parts)
                                table_text.append(row_text)
                        
                        # 只有当表格有内容时才添加
                        if table_text:
                            formatted_table = f"== 表格 {table_idx+1} ==\n" + "\n".join(table_text)
                            table_contents.append(formatted_table)
                
                # 组合段落和表格内容
                all_parts = []
                if paragraphs:
                    all_parts.append("\n".join(paragraphs))
                
                if table_contents:
                    all_parts.append("\n\n".join(table_contents))
                
                # 最终组合所有部分
                text_content = "\n\n".join(all_parts) if all_parts else "文档中没有找到文本内容"
            except Exception as e:
                text_content = f"无法解析Word文档: {str(e)}"
                
        elif file_type == 'excel':
            # 处理Excel文件
            import pandas as pd
            try:
                # 读取所有工作表
                excel_file = pd.ExcelFile(temp_path)
                sheet_texts = []
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    sheet_texts.append(f"== 工作表: {sheet_name} ==\n{df.to_string(index=False)}")
                
                text_content = "\n\n".join(sheet_texts)
            except Exception as e:
                text_content = f"无法解析Excel文件: {str(e)}"
                
        elif file_type == 'powerpoint':
            # 处理PowerPoint文件
            from pptx import Presentation
            try:
                prs = Presentation(temp_path)
                slide_texts = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                    
                    newline = '\n'
                    slide_texts.append(f"== 幻灯片 {slide_num} ==\n{newline.join(texts)}")
                
                text_content = "\n\n".join(slide_texts)
            except Exception as e:
                text_content = f"无法解析PowerPoint文件: {str(e)}"
        
        else:
            text_content = f"不支持的文件类型: {mime_type}"
        
        # 将提取的文本编码为base64
        text_base64 = base64.b64encode(text_content.encode('utf-8')).decode('utf-8')
        return text_base64, 'text/plain'
        
    except Exception as e:
        print(f"文档转换失败: {str(e)}")
        # 出错时返回原始文件
        return file_base64, mime_type
        
    finally:
        # 清理临时文件
        os.unlink(temp_path) 

def upload_file_to_gemini(file_source: str, api_key: str, file_mime_type: Optional[str] = None) -> tuple[str, str]:
    """
    将文件URL或base64数据上传至Gemini API
    
    Args:
        file_source: 文件URL或base64编码的文件数据
        api_key: Google API Key
        file_mime_type: 文件MIME类型，如果为None则自动检测
        
    Returns:
        tuple: (Gemini API返回的文件URI, 文件MIME类型)
    """
    file_data = None
    
    # 处理URL或base64格式的输入
    if file_source.startswith('http'):
        # 从URL下载文件
        response = requests.get(file_source)
        if response.status_code != 200:
            raise Exception(f"Failed to download file: {response.status_code}")
        file_data = response.content
        # 如果未指定MIME类型，尝试从响应头获取
        if not file_mime_type and 'Content-Type' in response.headers:
            file_mime_type = response.headers['Content-Type']
    elif file_source.startswith('data:'):
        # 处理base64格式
        mime_type, encoded_data = _get_mime_type_and_data(file_source)
        if mime_type:
            file_mime_type = mime_type
        file_data = base64.b64decode(encoded_data)
    else:
        # 假设是纯base64字符串
        try:
            file_data = base64.b64decode(file_source)
        except Exception:
            raise ValueError("Invalid file source format. Must be URL or base64 data.")
    
    # 默认MIME类型
    if not file_mime_type:
        file_mime_type = 'application/octet-stream'
    
    # 准备上传请求
    url = f"https://generativelanguage.googleapis.com/upload/v1beta/files?key={api_key}"
    
    # 创建临时文件名
    import uuid
    temp_filename = f"temp_file_{uuid.uuid4()}"
    
    # 准备multipart/form-data请求
    files = {
        'file': (temp_filename, file_data, file_mime_type)
    }
    
    # 发送上传请求
    upload_response = requests.post(url, files=files)
    
    if upload_response.status_code != 200:
        raise Exception(f"Failed to upload file: {upload_response.status_code}, {upload_response.text}")
    
    # 打印完整响应内容，帮助调试
    print(f"upload_response status: {upload_response.status_code}")
    print(f"upload_response headers: {upload_response.headers}")
    print(f"upload_response text: {upload_response.text}")
    
    # 检查响应内容是否为空
    if not upload_response.text.strip():
        raise Exception("File upload succeeded but response is empty")
    
    try:
        response_data = upload_response.json()
        if 'file' not in response_data or 'uri' not in response_data['file']:
            raise Exception(f"Invalid response format: {response_data}")
        
        # 返回文件URI
        return response_data['file']['uri'], file_mime_type
    except json.JSONDecodeError:
        # 如果无法解析JSON，记录错误并提供更多调试信息
        print(f"Failed to parse JSON response: {upload_response.text}")
        raise Exception(f"Failed to parse response as JSON. Raw response: {upload_response.text[:500]}")
